"""
Convert system-design-ppt.html (14 slides, fixed 1440x900) to PPTX.

Strategy: Playwright -> headless Chromium -> for each slide index,
toggle the .active class -> screenshot the #deck element -> embed each
PNG as a full-bleed image in a 16:10 PPTX.

Outputs:
  - docs/slides_png/slide-NN.png     (14 PNGs, 2x DPR for crispness)
  - docs/system-design-ppt.pptx       (14 slides, image-per-slide)
"""

import asyncio
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(r"d:\vsCode\python\blockchain\docs")
HTML = ROOT / "system-design-ppt.html"
PNG_DIR = ROOT / "slides_png"
PPTX_OUT = ROOT / "system-design-ppt.pptx"

WIDTH, HEIGHT = 1440, 900
N_SLIDES = 14
DPR = 2  # crisp output for projector / retina


async def render_slides():
    PNG_DIR.mkdir(parents=True, exist_ok=True)
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": WIDTH, "height": HEIGHT},
            device_scale_factor=DPR,
        )
        page = await ctx.new_page()
        url = HTML.resolve().as_uri()
        await page.goto(url, wait_until="networkidle")

        # 1. wait for fonts (Google Fonts CDN: Fraunces, IBM Plex, Noto Serif/Sans SC)
        await page.evaluate("() => document.fonts.ready")
        await page.wait_for_timeout(500)  # small extra buffer

        # 2. neutralize the runtime fit-scale transform; hide nav chrome; kill transitions
        await page.add_style_tag(content="""
            .slide { transition: none !important; }
            #nav, .nav { display: none !important; }
            html, body { background: #1F1C18 !important; }
        """)
        await page.evaluate("""
            () => {
                const deck = document.getElementById('deck');
                if (deck) {
                    deck.style.transform = 'none';
                    deck.style.boxShadow = 'none';
                }
            }
        """)

        for i in range(N_SLIDES):
            # 3. activate the i-th slide by toggling DOM class
            await page.evaluate(f"""
                () => {{
                    const slides = document.querySelectorAll('.slide');
                    slides.forEach(s => s.classList.remove('active', 'exit-left'));
                    slides[{i}].classList.add('active');
                    const prog = document.getElementById('progress');
                    if (prog) prog.style.width = (({i} + 1) / slides.length * 100) + '%';
                }}
            """)
            await page.wait_for_timeout(180)

            out = PNG_DIR / f"slide-{i + 1:02d}.png"
            deck = await page.query_selector("#deck")
            await deck.screenshot(path=str(out))
            print(f"  captured  slide-{i + 1:02d}.png")

        await browser.close()


def build_pptx():
    prs = Presentation()
    # 1440 x 900 = 16:10 -> 13.333" x 8.333"
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(8.333)
    blank = prs.slide_layouts[6]  # 'Blank' layout
    for i in range(N_SLIDES):
        slide = prs.slides.add_slide(blank)
        img = PNG_DIR / f"slide-{i + 1:02d}.png"
        slide.shapes.add_picture(
            str(img),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
    prs.save(str(PPTX_OUT))
    print(f"\n  PPTX saved: {PPTX_OUT}")


if __name__ == "__main__":
    print("[1/2] Rendering slides via Playwright/Chromium...")
    asyncio.run(render_slides())
    print("\n[2/2] Building PPTX...")
    build_pptx()
    print("\nDone.")
